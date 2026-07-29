"""발표 슬라이드 조립 검증. 네트워크 없이 폴백과 조립만 본다."""
import tempfile
import unittest
from pathlib import Path

from modules import ppt_report


class FallbackTests(unittest.TestCase):
    def test_chapter_headings_become_slides(self):
        text = 'I. 주제\n본문\nII. 방법\n본문\nIII. 결과'
        slides = ppt_report.fallback_slides('제목', text)
        self.assertEqual([s['heading'] for s in slides], ['주제', '방법', '결과'])

    def test_no_headings_still_gives_one_slide(self):
        slides = ppt_report.fallback_slides('발표', '장 제목이 없는 글')
        self.assertEqual(len(slides), 1)
        self.assertEqual(slides[0]['heading'], '발표')


class BuildTests(unittest.TestCase):
    def test_deck_round_trip(self):
        from pptx import Presentation
        slides = [
            {'heading': '탐구 질문', 'bullets': ['요점 하나', '요점 둘'],
             'notes': '말할 내용.'},
            {'heading': '빈 슬라이드', 'bullets': [], 'notes': ''},
        ]
        with tempfile.TemporaryDirectory() as tmp:
            out = Path(tmp) / 'deck.pptx'
            made = ppt_report.build_pptx('제목', '공통국어', slides, out)
            self.assertEqual(made, str(out))

            deck = Presentation(str(out))
            self.assertEqual(len(deck.slides), 3)   # 표지 + 2장
            self.assertEqual(deck.slides[1].shapes.title.text, '탐구 질문')
            self.assertIn('말할 내용',
                          deck.slides[1].notes_slide.notes_text_frame.text)
            # 빈 슬라이드에는 채울 자리 안내가 들어간다.
            texts = [shape.text_frame.text for shape in deck.slides[2].shapes
                     if shape.has_text_frame]
            self.assertTrue(any('채워' in text for text in texts))

    def test_long_bullets_are_clipped(self):
        from pptx import Presentation
        slides = [{'heading': 'h', 'bullets': ['가' * 200] * 10, 'notes': ''}]
        with tempfile.TemporaryDirectory() as tmp:
            out = Path(tmp) / 'deck.pptx'
            ppt_report.build_pptx('제목', '', slides, out)
            deck = Presentation(str(out))
            frame = deck.slides[1].placeholders[1].text_frame
            self.assertLessEqual(len(frame.paragraphs), ppt_report._MAX_BULLETS)
            self.assertLessEqual(len(frame.paragraphs[0].text),
                                 ppt_report._MAX_BULLET_CHARS)


if __name__ == '__main__':
    unittest.main()
