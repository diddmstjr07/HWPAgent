"""탐구 보고서를 발표 슬라이드(.pptx)로 만듭니다.

발표는 보고서와 다른 물건입니다. 문단을 옮겨 붙이면 글자만 가득한 슬라이드가
되므로, Codex가 장별 요점을 불릿으로 추리고 발표 메모(말할 내용)를 따로 답니다.
AI가 없거나 실패하면 보고서 제목·장 제목만으로라도 뼈대를 만들어 줍니다 —
학생이 내용을 채우는 편이 아무것도 없는 것보다 낫습니다.

원칙은 보고서와 같습니다. 슬라이드에 들어가는 내용은 보고서(=대화)에 있던 것만.
"""
import logging
import os
import re
from typing import Dict, List, Optional

import requests

logger = logging.getLogger(__name__)

_TIMEOUT = float(os.getenv('HWP_NODE_TIMEOUT_MS', '20000')) / 1000

# 한 슬라이드에 담을 수 있는 것에는 한계가 있다. 넘치면 잘라낸다.
_MAX_SLIDES = 12
_MAX_BULLETS = 6
_MAX_BULLET_CHARS = 60


def extract_report_text(path) -> str:
    """만들어진 HWP에서 본문 텍스트를 꺼냅니다(hwp-node 문서 모델 사용)."""
    base = (os.getenv('HWP_NODE_URL') or '').strip().rstrip('/')
    if not base:
        return ''
    headers = {}
    key = (os.getenv('HWP_NODE_API_KEY') or '').strip()
    if key:
        headers['X-API-Key'] = key

    session_id = None
    try:
        with requests.Session() as http:
            http.headers.update(headers)
            with open(path, 'rb') as handle:
                created = http.post(f'{base}/sessions', timeout=_TIMEOUT,
                                    files={'file': ('report.hwp', handle)})
            created.raise_for_status()
            session_id = created.json()['sessionId']
            document = http.get(f'{base}/sessions/{session_id}/document',
                                timeout=_TIMEOUT).json()
            return '\n'.join(
                (paragraph.get('text') or '')
                for section in document.get('sections') or []
                for paragraph in section.get('paragraphs') or [])
    except (requests.RequestException, ValueError, KeyError, OSError) as exc:
        logger.warning('[PPT] 보고서 텍스트 추출 실패: %s', type(exc).__name__)
        return ''
    finally:
        if session_id:
            try:
                requests.delete(f'{base}/sessions/{session_id}',
                                headers=headers, timeout=_TIMEOUT)
            except requests.RequestException:
                pass


_H1_RE = re.compile(r'^([IVX]+|[Ⅰ-Ⅻ])\.\s*(.+)$')


def fallback_slides(title: str, report_text: str) -> List[Dict]:
    """AI 없이 만드는 뼈대: 장 제목마다 슬라이드 하나. 내용은 학생이 채운다."""
    slides: List[Dict] = []
    for line in (report_text or '').splitlines():
        heading = _H1_RE.match(line.strip())
        if heading:
            slides.append({'heading': heading.group(2).strip(),
                           'bullets': [], 'notes': ''})
    return slides[:_MAX_SLIDES] or [{'heading': title or '탐구 발표',
                                     'bullets': [], 'notes': ''}]


def _clip(text: str, limit: int) -> str:
    text = str(text or '').strip()
    return text if len(text) <= limit else text[:limit - 1] + '…'


def build_pptx(title: str, subject: str, slides: List[Dict], out_path) -> Optional[str]:
    """슬라이드 목록을 .pptx로 조립합니다. 성공하면 경로, 실패하면 None."""
    try:
        from pptx import Presentation
        from pptx.util import Pt

        deck = Presentation()

        # 표지 — 제목 + 과목.
        cover = deck.slides.add_slide(deck.slide_layouts[0])
        cover.shapes.title.text = _clip(title or '탐구 발표', 80)
        if len(cover.placeholders) > 1:
            cover.placeholders[1].text = subject or ''

        for item in slides[:_MAX_SLIDES]:
            slide = deck.slides.add_slide(deck.slide_layouts[1])
            slide.shapes.title.text = _clip(item.get('heading') or '', 60)

            body = slide.placeholders[1].text_frame
            bullets = [b for b in (item.get('bullets') or []) if str(b).strip()]
            if bullets:
                for position, bullet in enumerate(bullets[:_MAX_BULLETS]):
                    paragraph = body.paragraphs[0] if position == 0 else body.add_paragraph()
                    paragraph.text = _clip(bullet, _MAX_BULLET_CHARS)
                    paragraph.font.size = Pt(20)
            else:
                # 뼈대만 있는 슬라이드. 채울 자리임을 밝혀 둔다.
                body.paragraphs[0].text = '(여기에 내용을 채워 주세요)'
                body.paragraphs[0].font.size = Pt(18)

            notes = str(item.get('notes') or '').strip()
            if notes:
                slide.notes_slide.notes_text_frame.text = notes

        deck.save(str(out_path))
        return str(out_path)
    except Exception as exc:  # noqa: BLE001 - 어떤 이유로든 못 만들면 없던 일로 한다
        logger.warning('[PPT] 조립 실패: %s: %s', type(exc).__name__, exc)
        return None
